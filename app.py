"""
JCC Assistant - Bible Study & Church Programs Chatbot
======================================================
Jubilee Celebration Center - AFM.
"""

import os
import re
import json
import base64
import urllib.parse
import urllib.request
from datetime import date
from pathlib import Path
from typing import List, Tuple, Optional

import gradio as gr
from supabase import create_client, Client
from openai import OpenAI
from docx import Document
from pptx import Presentation


# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
SUPABASE_URL = os.environ["SUPABASE_URL"]
SUPABASE_KEY = os.environ["SUPABASE_SERVICE_KEY"]
OPENAI_API_KEY = os.environ["OPENAI_API_KEY"]
ADMIN_PASSWORD = os.environ["ADMIN_PASSWORD"]

OPENAI_MODEL = "gpt-4o-mini"

sb: Client = create_client(SUPABASE_URL, SUPABASE_KEY)
oai = OpenAI(api_key=OPENAI_API_KEY)


# ---------------------------------------------------------------------------
# Logo as base64
# ---------------------------------------------------------------------------
LOGO_PATH = Path(__file__).parent / "jcc_logo.jpeg"
if LOGO_PATH.exists():
    with open(LOGO_PATH, "rb") as f:
        LOGO_B64 = base64.b64encode(f.read()).decode("ascii")
    LOGO_DATA_URI = f"data:image/jpeg;base64,{LOGO_B64}"
else:
    LOGO_DATA_URI = ""


# ---------------------------------------------------------------------------
# Document parsing
# ---------------------------------------------------------------------------
def parse_docx(file_path: str) -> str:
    doc = Document(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    sections = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        slide_lines.append(txt)
        if len(slide_lines) > 1:
            sections.append("\n".join(slide_lines))
    return "\n\n".join(sections)


def parse_study_document(file_path: str) -> str:
    lower = file_path.lower()
    if lower.endswith(".docx"):
        return parse_docx(file_path)
    elif lower.endswith(".pptx"):
        return parse_pptx(file_path)
    else:
        raise ValueError("Unsupported file type. Please upload .docx or .pptx.")


# ---------------------------------------------------------------------------
# Scripture lookup
# ---------------------------------------------------------------------------
# Recognize patterns like "John 3:16", "1 Corinthians 13:4-7", "Psalm 119:11"
SCRIPTURE_REGEX = re.compile(
    r"\b("
    r"(?:[123]\s*)?"               # optional book number like "1 ", "2 ", "3 "
    r"(?:[A-Z][a-zA-Z]+)"          # book name
    r"\s+\d+:\d+(?:[-\u2013]\d+)?" # chapter:verse(-verse)
    r")\b"
)


def lookup_scripture(reference: str, translation: str = "kjv") -> Optional[str]:
    """Fetch verse text from bible-api.com. Returns text or None on failure."""
    try:
        ref_clean = reference.replace("\u2013", "-").strip()
        url = f"https://bible-api.com/{urllib.parse.quote(ref_clean)}?translation={translation}"
        req = urllib.request.Request(url, headers={"User-Agent": "JCC-Assistant/1.0"})
        with urllib.request.urlopen(req, timeout=8) as r:
            data = json.loads(r.read().decode("utf-8"))
        text = (data.get("text") or "").strip()
        if text:
            return text
    except Exception as e:
        print(f"Scripture lookup failed for '{reference}': {e}")
    return None


def find_scriptures_in_text(text: str) -> List[str]:
    """Extract unique scripture references from text."""
    matches = SCRIPTURE_REGEX.findall(text)
    seen = []
    for m in matches:
        m_clean = re.sub(r"\s+", " ", m).strip()
        if m_clean and m_clean not in seen:
            seen.append(m_clean)
    return seen


# ---------------------------------------------------------------------------
# Suggested questions generator
# ---------------------------------------------------------------------------
SUGGEST_PROMPT = """Below is a Bible study document. Generate exactly 5 short, specific questions a group member might ask about THIS study after reading it.

Rules:
- Each question is 4-10 words.
- Questions must be answerable from the document content.
- Cover different sections / angles (not all the same topic).
- No generic Bible questions - they must be specific to this study.

Return a JSON array of 5 strings only. No preamble, no markdown, no code fences. Example:
["What does the study teach about X?", "What scriptures support Y?", ...]

STUDY DOCUMENT:
{document_text}
"""


def generate_suggested_questions(document_text: str) -> list:
    try:
        resp = oai.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[{"role": "user", "content": SUGGEST_PROMPT.format(document_text=document_text)}],
            temperature=0.5,
        )
        raw = resp.choices[0].message.content.strip()
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
            raw = raw.strip()
        questions = json.loads(raw)
        if isinstance(questions, list):
            return [str(q) for q in questions[:5]]
    except Exception as e:
        print(f"Question generation failed: {e}")
    return []


PROGRAMS_SUGGESTED = [
    "When is the next Couples Ministry event?",
    "What is the vision for Praise & Worship?",
    "What fundraising activities are planned for 2026?",
    "Who leads the Outreach Ministry?",
    "What events are happening this month?",
]


# ---------------------------------------------------------------------------
# Bible Study DB helpers
# ---------------------------------------------------------------------------
def list_bible_studies() -> List[Tuple[str, str]]:
    res = sb.table("bible_studies") \
        .select("id, week_of, title, presenter") \
        .order("week_of", desc=True) \
        .order("uploaded_at", desc=True) \
        .execute()
    options = []
    for row in res.data:
        label = f"{row['week_of']} - {row['title']}"
        if row.get("presenter"):
            label += f" ({row['presenter']})"
        options.append((label, row["id"]))
    return options


def get_bible_study(study_id: str) -> Optional[dict]:
    if not study_id:
        return None
    res = sb.table("bible_studies").select("*").eq("id", study_id).single().execute()
    return res.data


def upload_bible_study(file, title, presenter, week_of, password):
    if password != ADMIN_PASSWORD:
        return "❌ Incorrect admin password."
    if not file:
        return "❌ Please attach a .docx or .pptx file."
    if not title or not title.strip():
        return "❌ Title is required."
    if not week_of:
        return "❌ Week-of date is required."

    try:
        text = parse_study_document(file.name)
    except Exception as e:
        return f"❌ Could not parse the document: {e}"

    if len(text) < 50:
        return "❌ The document looks empty after parsing."

    questions = generate_suggested_questions(text)

    existing = sb.table("bible_studies").select("id").eq("week_of", week_of).execute()
    payload = {
        "title": title.strip(),
        "presenter": (presenter or "").strip() or None,
        "document_text": text,
        "suggested_questions": questions,
    }

    if existing.data:
        existing_id = existing.data[0]["id"]
        try:
            sb.table("bible_studies").update(payload).eq("id", existing_id).execute()
            return (
                f"✅ Replaced existing study for week of {week_of} with **{title}**.\n\n"
                f"Document length: {len(text):,} characters. "
                f"Generated {len(questions)} suggested questions."
            )
        except Exception as e:
            return f"❌ Database update failed: {e}"

    payload["week_of"] = week_of
    try:
        sb.table("bible_studies").insert(payload).execute()
    except Exception as e:
        return f"❌ Database insert failed: {e}"

    return (
        f"✅ Uploaded **{title}** for the week of {week_of}.\n\n"
        f"Document length: {len(text):,} characters. "
        f"Generated {len(questions)} suggested questions."
    )


def delete_bible_study(study_id, password):
    if password != ADMIN_PASSWORD:
        return "❌ Incorrect admin password.", gr.update()
    if not study_id:
        return "❌ Please select a study to delete.", gr.update()
    try:
        sb.table("bible_studies").delete().eq("id", study_id).execute()
    except Exception as e:
        return f"❌ Delete failed: {e}", gr.update()
    options = list_bible_studies()
    return "✅ Deleted.", gr.update(choices=options, value=options[0][1] if options else None)


# ---------------------------------------------------------------------------
# Church Programs context
# ---------------------------------------------------------------------------
def fetch_programs_context() -> str:
    ministries = sb.table("ministries").select("*").execute().data
    events = sb.table("events").select("*").order("event_date").execute().data
    notes = sb.table("ministry_notes").select("*").execute().data

    lines = ["# JCC 2026 MINISTRY PROGRAMS\n"]
    for m in ministries:
        lines.append(f"\n## {m['name']} Ministry")
        if m.get("lead"):
            lines.append(f"Led by: {m['lead']}")

        m_notes = [n for n in notes if n["ministry_id"] == m["id"]]
        for n in m_notes:
            heading = n["section"].replace("_", " ").title()
            lines.append(f"\n**{heading}:** {n['content']}")

        m_events = [e for e in events if e["ministry_id"] == m["id"]]
        if m_events:
            lines.append("\n**Events:**")
            for e in m_events:
                date_str = e.get("date_label") or e.get("event_date") or "TBD"
                line = f"- {date_str}: {e['title']}"
                if e.get("description"):
                    line += f" - {e['description']}"
                if e.get("format"):
                    line += f" ({e['format']})"
                lines.append(line)
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# System prompts (with follow-up + synonym handling)
# ---------------------------------------------------------------------------
BIBLE_STUDY_PROMPT = """You are the JCC Bible Study Assistant for Jubilee Celebration Center - AFM.

Today's date is {today}.

You answer questions about the Bible study document below. Your answers MUST come only from the document - do not add interpretation, outside scripture, or commentary.

RULES:
- Quote directly from the document using quotation marks for the key passage.
- Cite the section heading, number, or slide where the quote appears.
- Treat follow-up questions as connected to the prior exchange. If the user previously asked about "main points" and follows up with "tell me more about [topic that appeared in your answer]", that topic IS in the document - dig into the relevant section and elaborate. Do NOT decline.
- Match user phrasing to document content semantically. If they ask about a concept the document covers under a different word (e.g., they say "what is Logos?" and the doc has a section "Logos and Rhema"), treat that as a valid match and answer from that section.
- Only respond with "This isn't covered in this week's study" when the topic GENUINELY isn't in the document at all.
- Do not speculate or fill gaps with general Bible knowledge.
- Keep answers focused.

---
THIS WEEK'S BIBLE STUDY
Title: {title}
Presenter: {presenter}
Week of: {week_of}

{document_text}
---
"""

PROGRAMS_PROMPT = """You are the JCC Programs Assistant for Jubilee Celebration Center - AFM.

Today's date is {today}. Use this when answering questions about "next", "upcoming", "past", or "today's" events.

You answer questions about JCC's 2026 ministry programs, events, leads, goals, and activities from the data below.

SYNONYM AND SEMANTIC MAPPING (very important):
- "Evangelism" / "missions" / "community service" → Outreach Ministry
- "Music" / "worship team" / "choir" → Praise & Worship Ministry
- "Men's group" / "brothers" → Men of God Ministry
- "Ladies" / "women's group" / "sisters" → Women of God Ministry
- "Marriage" / "married couples" → Couples Ministry
- "Greeters" / "ushers" / "welcome team" → Hospitality Ministry
- "Building" / "repairs" / "facilities" → Maintenance & Repair Ministry
- "Money" / "donations" / "giving" → Fundraising Ministry
Always recognize these synonyms and answer from the corresponding ministry's data.

GUIDANCE:
- Be specific. When asked about events, give the date, ministry, and format.
- For "next" or "upcoming" events, only consider events with dates AFTER today ({today}). If the next chronological event is in the past, skip it and give the next future one. If all are in the past, say there are no upcoming events for that ministry.
- For questions about people, look across ALL ministries - leads are listed at each ministry section. Example: "Who is the Pastor?" → Pastor Tabu Bere leads the Outreach Ministry.
- If asked about "Bible study" or a "Bible Study Ministry": there is no separate Bible Study Ministry - Bible studies are presented rotationally by groups within the congregation. Tell the user this and suggest they switch to "Bible Study" mode for content questions.
- Only say "I don't have that information" when the topic GENUINELY isn't anywhere in the data, even by synonym.
- Do not invent events, dates, or leads.

---
JCC 2026 PROGRAMS DATA
{programs_context}
---
"""


# ---------------------------------------------------------------------------
# Chat
# ---------------------------------------------------------------------------
def chat(message, history, mode, study_id):
    if not message or not message.strip():
        return "Please type a question."

    today_str = date.today().isoformat()

    if mode == "Bible Study":
        study = get_bible_study(study_id) if study_id else None
        if not study:
            return (
                "Please select a Bible study from the dropdown first. "
                "If the dropdown is empty, click Refresh or ask an admin to upload one."
            )
        system_prompt = BIBLE_STUDY_PROMPT.format(
            today=today_str,
            title=study["title"],
            presenter=study.get("presenter") or "Not specified",
            week_of=study["week_of"],
            document_text=study["document_text"],
        )
    else:
        try:
            context = fetch_programs_context()
        except Exception as e:
            return f"Could not load programs data: {e}"
        system_prompt = PROGRAMS_PROMPT.format(today=today_str, programs_context=context)

    messages = [{"role": "system", "content": system_prompt}]
    for h in history:
        messages.append({"role": h["role"], "content": h["content"]})
    messages.append({"role": "user", "content": message})

    try:
        resp = oai.chat.completions.create(
            model=OPENAI_MODEL,
            messages=messages,
            temperature=0.2,
        )
        return resp.choices[0].message.content
    except Exception as e:
        return f"OpenAI request failed: {e}"


# ---------------------------------------------------------------------------
# Dynamic suggested questions
# ---------------------------------------------------------------------------
def get_suggestions(mode: str, study_id: str) -> List[str]:
    if mode == "Church Programs":
        return PROGRAMS_SUGGESTED
    if not study_id:
        return ["(Select a study to see suggestions)"]
    try:
        study = get_bible_study(study_id)
        if study and study.get("suggested_questions"):
            qs = study["suggested_questions"]
            if isinstance(qs, list) and qs:
                return qs
        return [
            "What were the main points?",
            "What scriptures are referenced?",
            "Summarize the conclusion",
            "What does the study say about the key topic?",
            "What is the application for our lives?",
        ]
    except Exception as e:
        print(f"get_suggestions failed: {e}")
        return []


def get_study_text(mode: str, study_id: str) -> str:
    """Return the document text for the currently selected study (for the side panel)."""
    if mode != "Bible Study" or not study_id:
        return ""
    try:
        study = get_bible_study(study_id)
        if not study:
            return ""
        title = study["title"]
        presenter = study.get("presenter") or ""
        week_of = study["week_of"]
        header = f"### {title}\n\n**Week of:** {week_of}"
        if presenter:
            header += f"   ·   **Presenter:** {presenter}"
        header += "\n\n---\n\n"
        return header + study["document_text"]
    except Exception as e:
        return f"Error loading study text: {e}"


def update_after_change(mode, study_id):
    """Update suggestions AND study text after mode or study changes."""
    suggestions = get_suggestions(mode, study_id)
    updates = []
    for i in range(5):
        if i < len(suggestions):
            updates.append(gr.update(value=suggestions[i], visible=True))
        else:
            updates.append(gr.update(visible=False))
    study_text = get_study_text(mode, study_id)
    text_visible = bool(study_text)
    return (*updates, gr.update(value=study_text, visible=text_visible))


def refresh_studies_full(mode):
    options = list_bible_studies()
    if not options:
        suggestions = get_suggestions(mode, None)
        sugg_updates = []
        for i in range(5):
            if i < len(suggestions):
                sugg_updates.append(gr.update(value=suggestions[i], visible=True))
            else:
                sugg_updates.append(gr.update(visible=False))
        return (
            gr.update(choices=[], value=None),
            *sugg_updates,
            gr.update(value="", visible=False),
        )
    new_value = options[0][1]
    suggestions = get_suggestions(mode, new_value)
    sugg_updates = []
    for i in range(5):
        if i < len(suggestions):
            sugg_updates.append(gr.update(value=suggestions[i], visible=True))
        else:
            sugg_updates.append(gr.update(visible=False))
    study_text = get_study_text(mode, new_value)
    return (
        gr.update(choices=options, value=new_value),
        *sugg_updates,
        gr.update(value=study_text, visible=bool(study_text)),
    )


# ---------------------------------------------------------------------------
# Scripture lookup handler (Tier-1 fix: clicking refs in the chat is hard
# in stock Gradio; we add a separate "look up a verse" input instead).
# ---------------------------------------------------------------------------
def lookup_scripture_ui(ref, translation):
    if not ref or not ref.strip():
        return "Type a verse reference like *John 3:16* or *Romans 10:17*."
    text = lookup_scripture(ref.strip(), translation or "kjv")
    if text is None:
        return f"Could not find **{ref}** in the {translation.upper()} translation. Check the spelling and try again."
    return f"**{ref}** ({translation.upper()})\n\n> {text}"


# ---------------------------------------------------------------------------
# UI
# ---------------------------------------------------------------------------
CUSTOM_CSS = """
.gradio-container {
    font-family: 'Inter', 'Helvetica Neue', system-ui, sans-serif !important;
    max-width: 1400px !important;
    margin: 0 auto !important;
}
#jcc-hero {
    background: linear-gradient(135deg, #1B2A4E 0%, #2C4170 100%);
    border-radius: 16px;
    padding: 28px 32px;
    margin-bottom: 18px;
    color: white;
    display: flex;
    align-items: center;
    gap: 24px;
    box-shadow: 0 8px 24px rgba(27, 42, 78, 0.18);
    position: relative;
    overflow: hidden;
}
#jcc-hero::after {
    content: "";
    position: absolute;
    bottom: 0; left: 0; right: 0;
    height: 4px;
    background: linear-gradient(90deg, #C9A55C 0%, #E4CC8E 50%, #C9A55C 100%);
}
#jcc-hero img.logo {
    width: 88px;
    height: 88px;
    border-radius: 50%;
    background: white;
    padding: 6px;
    box-shadow: 0 4px 12px rgba(0,0,0,0.2);
    flex-shrink: 0;
}
#jcc-hero .titles h1 {
    font-size: 1.9em !important;
    font-weight: 700 !important;
    margin: 0 0 4px 0 !important;
    color: white !important;
    letter-spacing: -0.5px;
}
#jcc-hero .titles .church-name {
    font-size: 0.85em;
    color: #C9A55C;
    letter-spacing: 3px;
    font-weight: 600;
    margin-bottom: 6px;
    text-transform: uppercase;
}
#jcc-hero .titles .tagline {
    font-size: 0.95em;
    color: #cbd5e1;
    margin: 0;
}
.sidebar-card {
    background: white;
    border-radius: 12px;
    padding: 16px;
    border: 1px solid #e5e7eb;
    margin-bottom: 12px;
}
.sidebar-card h3 {
    color: #1B2A4E;
    font-size: 0.78em;
    letter-spacing: 1.5px;
    text-transform: uppercase;
    margin: 0 0 12px 0;
    font-weight: 700;
    border-left: 3px solid #C9A55C;
    padding-left: 10px;
}
.suggest-btn button {
    background: white !important;
    border: 1px solid #e5e7eb !important;
    color: #1B2A4E !important;
    text-align: left !important;
    font-weight: 500 !important;
    font-size: 0.88em !important;
    padding: 10px 12px !important;
    line-height: 1.35 !important;
    white-space: normal !important;
    height: auto !important;
    min-height: 40px !important;
    transition: all 0.15s ease;
    width: 100% !important;
    justify-content: flex-start !important;
}
.suggest-btn button:hover {
    background: #1B2A4E !important;
    color: white !important;
    border-color: #1B2A4E !important;
    transform: translateX(2px);
}
.tab-nav button { font-weight: 500 !important; }
.tab-nav button.selected {
    color: #1B2A4E !important;
    border-bottom-color: #C9A55C !important;
}
#study-text-panel {
    background: #fefcf7;
    border: 1px solid #e8dfc7;
    border-radius: 12px;
    padding: 18px 20px;
    max-height: 560px;
    overflow-y: auto;
    font-size: 0.92em;
    line-height: 1.55;
    color: #1f2937;
}
#study-text-panel h3 { color: #1B2A4E; margin-top: 0; }
footer { display: none !important; }
"""


theme = gr.themes.Soft(
    primary_hue=gr.themes.colors.slate,
    secondary_hue=gr.themes.colors.amber,
    neutral_hue=gr.themes.colors.slate,
    font=[gr.themes.GoogleFont("Inter"), "system-ui", "sans-serif"],
).set(
    button_primary_background_fill="#1B2A4E",
    button_primary_background_fill_hover="#0F1A35",
    button_primary_text_color="white",
    body_background_fill="#F7F3EC",
    block_background_fill="white",
    block_border_color="#e5e7eb",
)


with gr.Blocks(title="JCC Assistant", theme=theme, css=CUSTOM_CSS) as demo:

    logo_img_html = (
        f'<img class="logo" src="{LOGO_DATA_URI}" alt="JCC Logo"/>'
        if LOGO_DATA_URI else ""
    )
    gr.HTML(f"""
    <div id="jcc-hero">
        {logo_img_html}
        <div class="titles">
            <div class="church-name">Jubilee Celebration Center &mdash; AFM</div>
            <h1>JCC Assistant</h1>
            <p class="tagline">Ask about this week's Bible study or church programs and events for 2026.</p>
        </div>
    </div>
    """)

    with gr.Tabs():
        # ============================================================
        # CHAT TAB
        # ============================================================
        with gr.Tab("Chat"):
            with gr.Row():
                # LEFT SIDEBAR
                with gr.Column(scale=1, min_width=240):
                    with gr.Group(elem_classes=["sidebar-card"]):
                        gr.HTML("<h3>Mode</h3>")
                        mode = gr.Radio(
                            choices=["Bible Study", "Church Programs"],
                            value="Bible Study",
                            show_label=False,
                            container=False,
                        )

                    with gr.Group(elem_classes=["sidebar-card"]):
                        gr.HTML("<h3>Bible Study</h3>")
                        study_dropdown = gr.Dropdown(
                            choices=list_bible_studies(),
                            show_label=False,
                            container=False,
                        )
                        refresh_btn = gr.Button("Refresh list", size="sm")

                    with gr.Group(elem_classes=["sidebar-card"]):
                        gr.HTML("<h3>Suggested Questions</h3>")
                        suggest_btns = [
                            gr.Button("", elem_classes=["suggest-btn"], visible=False)
                            for _ in range(5)
                        ]

                # MAIN CHAT
                with gr.Column(scale=2):
                    chatbot = gr.Chatbot(
                        type="messages",
                        height=560,
                        avatar_images=(None, str(LOGO_PATH) if LOGO_PATH.exists() else None),
                        show_label=False,
                        show_copy_button=True,
                    )
                    msg = gr.Textbox(
                        placeholder="Ask a question…",
                        show_label=False,
                        container=False,
                        autofocus=True,
                    )

                # RIGHT PANEL: study document text
                with gr.Column(scale=2):
                    with gr.Group(elem_classes=["sidebar-card"]):
                        gr.HTML("<h3>Study Document</h3>")
                        study_text_display = gr.Markdown(
                            value="",
                            elem_id="study-text-panel",
                            visible=False,
                        )

            # ---------- chat plumbing ----------
            def respond(message, history, mode_val, study_id):
                if not message or not message.strip():
                    return "", history
                history = history + [{"role": "user", "content": message}]
                reply = chat(message, history[:-1], mode_val, study_id)
                history = history + [{"role": "assistant", "content": reply}]
                return "", history

            msg.submit(
                respond,
                inputs=[msg, chatbot, mode, study_dropdown],
                outputs=[msg, chatbot],
            )

            def send_suggested(question, history, mode_val, study_id):
                if not question or not question.strip():
                    return history
                history = history + [{"role": "user", "content": question}]
                reply = chat(question, history[:-1], mode_val, study_id)
                history = history + [{"role": "assistant", "content": reply}]
                return history

            for btn in suggest_btns:
                btn.click(
                    fn=send_suggested,
                    inputs=[btn, chatbot, mode, study_dropdown],
                    outputs=[chatbot],
                )

            # Mode/study changes update suggestions AND study text
            mode.change(
                fn=update_after_change,
                inputs=[mode, study_dropdown],
                outputs=[*suggest_btns, study_text_display],
            )
            study_dropdown.change(
                fn=update_after_change,
                inputs=[mode, study_dropdown],
                outputs=[*suggest_btns, study_text_display],
            )
            refresh_btn.click(
                fn=refresh_studies_full,
                inputs=[mode],
                outputs=[study_dropdown, *suggest_btns, study_text_display],
            )

            # Init on load
            demo.load(
                fn=update_after_change,
                inputs=[mode, study_dropdown],
                outputs=[*suggest_btns, study_text_display],
            )

        # ============================================================
        # SCRIPTURE LOOKUP TAB
        # ============================================================
        with gr.Tab("Scripture Lookup"):
            gr.Markdown(
                "### Look up a Bible verse\n"
                "Type a reference like `John 3:16`, `Romans 10:17`, `1 Corinthians 13:4-7`, or `Psalm 119:11`."
            )
            with gr.Row():
                ref_input = gr.Textbox(
                    label="Verse reference",
                    placeholder="e.g. John 1:1-3",
                    scale=3,
                )
                translation_choice = gr.Dropdown(
                    label="Translation",
                    choices=[
                        ("King James Version (KJV)", "kjv"),
                        ("World English Bible (WEB)", "web"),
                        ("American Standard Version (ASV)", "asv"),
                        ("Berean Standard Bible (BSB)", "bsb"),
                    ],
                    value="kjv",
                    scale=1,
                )
            lookup_btn = gr.Button("Look up verse", variant="primary")
            verse_output = gr.Markdown()

            lookup_btn.click(
                fn=lookup_scripture_ui,
                inputs=[ref_input, translation_choice],
                outputs=verse_output,
            )
            ref_input.submit(
                fn=lookup_scripture_ui,
                inputs=[ref_input, translation_choice],
                outputs=verse_output,
            )

            gr.Markdown(
                "---\n*Powered by bible-api.com. NIV/NKJV translations require a paid API "
                "and are not available here. KJV, WEB, ASV, and BSB are free.*"
            )

        # ============================================================
        # ADMIN TAB
        # ============================================================
        with gr.Tab("Admin"):
            gr.Markdown(
                "### Upload a Bible Study\n"
                "Accepts `.docx` or `.pptx` files. "
                "If a study already exists for the same week-of date, it will be **replaced**. "
                "Five suggested questions are auto-generated for each upload."
            )

            with gr.Row():
                with gr.Column(scale=2):
                    up_title = gr.Textbox(
                        label="Title",
                        placeholder="e.g. The Foundation of the Word",
                    )
                    up_presenter = gr.Textbox(
                        label="Presenter",
                        placeholder="e.g. Elder, Group 3",
                    )
                    up_week = gr.Textbox(
                        label="Week of (YYYY-MM-DD)",
                        value=str(date.today()),
                    )
                with gr.Column(scale=1):
                    up_password = gr.Textbox(
                        label="Admin Password",
                        type="password",
                    )
                    up_file = gr.File(
                        label="Document (.docx / .pptx)",
                        file_types=[".docx", ".pptx"],
                    )

            up_button = gr.Button("Upload Study", variant="primary")
            up_status = gr.Markdown()

            up_button.click(
                fn=upload_bible_study,
                inputs=[up_file, up_title, up_presenter, up_week, up_password],
                outputs=up_status,
            )

            gr.Markdown("---\n### Delete a Bible Study")
            with gr.Row():
                del_dropdown = gr.Dropdown(
                    label="Select study to delete",
                    choices=list_bible_studies(),
                    scale=4,
                )
                del_refresh = gr.Button("Refresh list", scale=1, size="sm")
                del_password = gr.Textbox(
                    label="Admin Password",
                    type="password",
                    scale=2,
                )
            del_button = gr.Button("Delete Selected", variant="stop")
            del_status = gr.Markdown()

            def refresh_del():
                options = list_bible_studies()
                return gr.update(choices=options, value=options[0][1] if options else None)

            del_refresh.click(fn=refresh_del, outputs=del_dropdown)
            del_button.click(
                fn=delete_bible_study,
                inputs=[del_dropdown, del_password],
                outputs=[del_status, del_dropdown],
            )

    gr.HTML(
        "<div style='text-align:center; color:#9ca3af; font-size:0.85em; padding:12px;'>"
        "JCC Assistant - Prototype. The bot only answers from loaded study and programs data."
        "</div>"
    )


if __name__ == "__main__":
    demo.launch()